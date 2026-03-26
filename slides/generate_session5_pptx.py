"""Generate session5_files_errors_project.pptx from the Session 5 course content.

Usage:
    pip install python-pptx
    python slides/generate_session5_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)   # dark title backgrounds
TEAL   = RGBColor(0x00, 0x70, 0xC0)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # subtle slide background
GREEN  = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xC5, 0x56, 0x12)
GRAY   = RGBColor(0x44, 0x44, 0x44)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)  # VS-Code dark background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, bold=False, italic=False,
                color=GRAY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    """Full-width navy title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.4))
    fill_solid(bar, NAVY)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xF5)


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Dark VS-Code-style code block."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    fill_solid(box, CODE_BG)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)


def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # completely blank layout

    # ── Slide 1: Title / Welcome ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
                "📁 Files, Errors & Your First Project", font_size=44, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
                "Session 5: File I/O, Error Handling & Mini-Project",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(3), Inches(3.9), Inches(7.33), Pt(3))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                "Duration: 1 hour 40 minutes",
                font_size=20, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.5),
                "Build real programs that read files and handle errors! 🐍",
                font_size=18, italic=True, color=RGBColor(0xA0, 0xC8, 0xFF),
                align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📋 What We'll Cover Today")

    agenda = [
        "Recap & Assignment 4 Review",
        "Reading files: open(), read(), readlines(), with",
        "Writing & Appending files",
        "Error handling: try, except, finally",
        "Common Exceptions",
        "f-strings deep dive",
        "🎓 Course recap",
        "🛠️ Mini-Project: Command-Line To-Do App",
    ]

    txb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(agenda, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}.  {item}"
        run.font.size = Pt(22)
        run.font.color.rgb = NAVY if i % 2 == 1 else TEAL

    # ── Slide 3: Recap & Assignment 4 Review ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🧠 Recap & Assignment 4 Review")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Quick check — what's the output?",
                font_size=20, bold=True, color=NAVY)

    add_code_block(s,
                   "# What's the output?\n"
                   'data = {"a": 1, "b": 2, "c": 3}\n'
                   "result = [k for k, v in data.items() if v > 1]\n"
                   "print(result)\n\n"
                   "# Fix the bug:\n"
                   'students = [{"name": "Alice", "grade": 90}]\n'
                   "for s in students:\n"
                   "    print(s[name])   # Error — what's wrong?\n\n"
                   "# What does this create?\n"
                   "matrix = [[i * j for j in range(1, 4)] for i in range(1, 4)]",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.8),
                   font_size=14)

    # ── Slide 4: Reading Files ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📂 Reading Files")

    add_code_block(s,
                   '# Basic file reading\n'
                   'file = open("data.txt", "r")   # "r" = read mode\n'
                   "content = file.read()          # Read entire file as string\n"
                   "print(content)\n"
                   "file.close()                   # Always close!\n\n"
                   "# Better way — with statement (auto-closes)\n"
                   'with open("data.txt", "r") as file:\n'
                   "    content = file.read()\n"
                   "# File is automatically closed here\n\n"
                   "# Read line by line\n"
                   'with open("data.txt", "r") as file:\n'
                   "    lines = file.readlines()   # Returns list of strings\n"
                   "    for line in lines:\n"
                   "        print(line.strip())    # .strip() removes \\n",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=14)

    # ── Slide 5: File Modes ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📂 File Modes")

    modes_data = [
        ('"r"',  "Read (default) — file must exist"),
        ('"w"',  "Write — creates new or overwrites"),
        ('"a"',  "Append — adds to end of file"),
        ('"r+"', "Read and write"),
        ('"rb"', "Read binary (images, etc.)"),
    ]

    tbl = s.shapes.add_table(
        len(modes_data) + 1, 2,
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(3.0),
    ).table
    tbl.columns[0].width = Inches(1.4)
    tbl.columns[1].width = Inches(4.4)

    for ci, hdr in enumerate(["Mode", "Description"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE

    for ri, (mode, desc) in enumerate(modes_data, 1):
        bg = LIGHT if ri % 2 == 1 else WHITE
        for ci, txt in enumerate([mode, desc]):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = txt
            run.font.size = Pt(14)
            run.font.color.rgb = TEAL if ci == 0 else GRAY

    add_code_block(s,
                   "# Reading line by line (memory efficient for large files)\n"
                   'with open("big_file.txt", "r") as file:\n'
                   "    for line in file:           # Iterate directly!\n"
                   "        print(line.strip())\n\n"
                   "# Read and process each line\n"
                   'with open("names.txt", "r") as file:\n'
                   "    names = [line.strip() for line in file]\n"
                   "print(names)",
                   Inches(6.9), Inches(1.6), Inches(6.0), Inches(3.4),
                   font_size=13)

    # ── Slide 6: Writing & Appending Files ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "✏️ Writing & Appending Files")

    add_code_block(s,
                   "# Write mode — creates new file (or overwrites!)\n"
                   'with open("output.txt", "w") as file:\n'
                   '    file.write("Hello, File!\\n")\n'
                   '    file.write("Second line\\n")\n\n'
                   "# Write multiple lines at once\n"
                   'lines = ["Line 1\\n", "Line 2\\n", "Line 3\\n"]\n'
                   'with open("output.txt", "w") as file:\n'
                   "    file.writelines(lines)\n\n"
                   "# Append mode — adds to existing file\n"
                   'with open("log.txt", "a") as file:\n'
                   '    file.write("New log entry\\n")\n\n'
                   "# Write structured data\n"
                   'tasks = ["Buy milk", "Call mom", "Study Python"]\n'
                   'with open("tasks.txt", "w") as file:\n'
                   "    for task in tasks:\n"
                   '        file.write(task + "\\n")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=14)

    # ── Slide 7: Error Handling: try/except ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛡️ Error Handling: try / except")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Python raises exceptions when something goes wrong.",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# Without error handling — program crashes!\n"
                   'age = int(input("Enter age: "))  # What if user types "abc"?\n\n'
                   "# With error handling — graceful!\n"
                   "try:\n"
                   '    age = int(input("Enter age: "))\n'
                   "    print(f\"You are {age} years old.\")\n"
                   "except ValueError:\n"
                   '    print("Please enter a valid number!")\n\n'
                   "# The flow:\n"
                   "# try block runs\n"
                   "# If an error occurs  →  jump to except\n"
                   "# If no error         →  skip except\n"
                   "# Program continues either way",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.8),
                   font_size=14)

    # ── Slide 8: More Error Handling Patterns ────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛡️ More Error Handling Patterns")

    add_code_block(s,
                   "# Catch multiple exception types\n"
                   "try:\n"
                   '    result = int(input("Number: ")) / int(input("Divide by: "))\n'
                   '    print(f"Result: {result}")\n'
                   "except ValueError:\n"
                   '    print("Please enter valid numbers!")\n'
                   "except ZeroDivisionError:\n"
                   '    print("Cannot divide by zero!")\n\n'
                   "# finally — always runs (great for cleanup)\n"
                   "try:\n"
                   '    file = open("data.txt", "r")\n'
                   "    content = file.read()\n"
                   "except FileNotFoundError:\n"
                   '    print("File not found!")\n'
                   '    content = ""\n'
                   "finally:\n"
                   '    print("Done attempting to read file.")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=14)

    # ── Slide 9: Common Exceptions ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "⚠️ Common Exceptions")

    exc_data = [
        ("ValueError",        'Wrong type: int("abc")'),
        ("ZeroDivisionError", "Dividing by zero: 10 / 0"),
        ("FileNotFoundError", "File doesn't exist"),
        ("IndexError",        "List index out of range: lst[99]"),
        ("KeyError",          'Dict key missing: d["missing"]'),
        ("TypeError",         'Wrong operation: "hi" + 5'),
        ("NameError",         "Variable not defined: print(x)"),
    ]

    tbl = s.shapes.add_table(
        len(exc_data) + 1, 2,
        Inches(0.8), Inches(1.6), Inches(8.5), Inches(3.8),
    ).table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(5.5)

    for ci, hdr in enumerate(["Exception", "When it occurs"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE

    for ri, (exc, when) in enumerate(exc_data, 1):
        bg = LIGHT if ri % 2 == 1 else WHITE
        for ci, txt in enumerate([exc, when]):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = txt
            run.font.size = Pt(13)
            run.font.color.rgb = TEAL if ci == 0 else GRAY

    add_code_block(s,
                   "# Catching any exception (use sparingly!)\n"
                   "try:\n"
                   "    risky_operation()\n"
                   "except Exception as e:\n"
                   '    print(f"An error occurred: {e}")',
                   Inches(0.8), Inches(5.65), Inches(11.5), Inches(1.55),
                   font_size=14)

    # ── Slide 10: f-strings Deep Dive ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎨 f-strings Deep Dive")

    add_code_block(s,
                   'name    = "Alice"\n'
                   "score   = 95.678\n"
                   "balance = 1234567.89\n\n"
                   "# Basic\n"
                   'print(f"Hello, {name}!")\n\n'
                   "# Expressions inside {}\n"
                   'print(f"Score: {score:.2f}")       # 95.68 (2 decimal places)\n'
                   'print(f"Score: {score:.0f}")       # 96   (no decimals)\n'
                   'print(f"Name: {name:>10}")         # right-align in 10 chars\n'
                   'print(f"Name: {name:<10}|")        # left-align\n'
                   'print(f"Name: {name:^10}")         # center\n\n'
                   "# Numbers\n"
                   'print(f"Balance: {balance:,.2f}")  # 1,234,567.89 (with commas)\n'
                   'print(f"Pi: {3.14159:.3f}")        # 3.142\n\n'
                   "# Expressions\n"
                   'print(f"Double: {score * 2}")\n'
                   'print(f"Upper: {name.upper()}")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=13)

    # ── Slide 11: Course Recap ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎓 Course Recap — Everything You've Learned!")

    add_code_block(s,
                   "Session 1  ──►  print(), variables, data types, input(), operators\n"
                   "                          │\n"
                   "Session 2  ──►  if/elif/else, while, for, range, break, continue\n"
                   "                          │\n"
                   "Session 3  ──►  def, parameters, return, scope, modules\n"
                   "                          │\n"
                   "Session 4  ──►  lists, tuples, dicts, comprehensions\n"
                   "                          │\n"
                   "Session 5  ──►  file I/O, try/except, f-strings",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.2),
                   font_size=15)

    add_textbox(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.45),
                "You can now build real programs that:",
                font_size=20, bold=True, color=NAVY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "Take user input and make decisions",
        "Repeat actions efficiently",
        "Organize code into reusable functions",
        "Work with complex data",
        "Read/write files and handle errors",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  ✅  {item}"
        run.font.size = Pt(18)
        run.font.color.rgb = GRAY

    # ── Slide 12: Mini-Project Design ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛠️ Mini-Project: To-Do App — Step 1: Design")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "We'll build a Command-Line To-Do App together!",
                font_size=20, bold=True, color=NAVY)

    add_textbox(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.45),
                "Features:", font_size=18, bold=True, color=TEAL)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, feat in enumerate([
        "Add tasks",
        "View all tasks",
        "Mark task as complete",
        "Delete a task",
        "Save to file (tasks persist!)",
        "Load on startup",
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {feat}"
        run.font.size = Pt(18)
        run.font.color.rgb = GRAY

    add_code_block(s,
                   "=== To-Do App ===\n"
                   "1. Add task\n"
                   "2. View tasks\n"
                   "3. Mark complete\n"
                   "4. Delete task\n"
                   "5. Quit\n\n"
                   "Choose an option: _",
                   Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.2),
                   font_size=16)

    # ── Slide 13: Mini-Project Step 2: Data Structure ─────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛠️ Mini-Project — Step 2: Data Structure")

    add_code_block(s,
                   'FILENAME = "tasks.txt"\n\n'
                   "# Each task is a dict:\n"
                   "task = {\n"
                   '    "description": "Buy groceries",\n'
                   '    "done": False\n'
                   "}\n\n"
                   "# All tasks in a list:\n"
                   "tasks = [\n"
                   '    {"description": "Buy groceries", "done": False},\n'
                   '    {"description": "Study Python",  "done": True},\n'
                   "]",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2),
                   font_size=16)

    # ── Slide 14: Mini-Project Step 3: Core Functions ─────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛠️ Mini-Project — Step 3: Core Functions")

    add_code_block(s,
                   "def load_tasks():\n"
                   "    tasks = []\n"
                   "    try:\n"
                   "        with open(FILENAME, \"r\") as f:\n"
                   "            for line in f:\n"
                   "                parts = line.strip().split(\"|\")\n"
                   "                tasks.append({\n"
                   '                    "description": parts[0],\n'
                   '                    "done": parts[1] == "True"\n'
                   "                })\n"
                   "    except FileNotFoundError:\n"
                   "        pass   # No file yet — start fresh\n"
                   "    return tasks\n\n"
                   "def save_tasks(tasks):\n"
                   "    with open(FILENAME, \"w\") as f:\n"
                   "        for task in tasks:\n"
                   "            f.write(f\"{task['description']}|{task['done']}\\n\")",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=14)

    # ── Slide 15: Mini-Project Step 4: UI Functions ───────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛠️ Mini-Project — Step 4: UI Functions")

    add_code_block(s,
                   "def view_tasks(tasks):\n"
                   "    if not tasks:\n"
                   '        print("No tasks yet!")\n'
                   "        return\n"
                   '    print("\\n=== Your Tasks ===")\n'
                   "    for i, task in enumerate(tasks, 1):\n"
                   '        status = "✓" if task["done"] else "○"\n'
                   '        print(f"{i}. [{status}] {task[\'description\']}")\n\n'
                   "def add_task(tasks):\n"
                   '    desc = input("Task description: ").strip()\n'
                   "    if desc:\n"
                   '        tasks.append({"description": desc, "done": False})\n'
                   "        print(f\"Added: '{desc}'\")\n"
                   "    else:\n"
                   '        print("Task cannot be empty!")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=14)

    # ── Slide 16: Practice Exercises ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎯 Practice Exercises")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Try these on your own:", font_size=22, bold=True, color=NAVY)

    exercises = [
        ("Exercise 1",
         "Write a function that reads a file and counts the number of lines,\n"
         "  words, and characters."),
        ("Exercise 2",
         "Write a try/except block that:\n"
         "  • Opens a file the user specifies\n"
         "  • Handles FileNotFoundError gracefully\n"
         "  • Prints the first 5 lines"),
        ("Exercise 3",
         "Extend the To-Do App to support a \"search tasks\" feature."),
    ]

    y = Inches(2.2)
    for title, body in exercises:
        add_textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.4),
                    title, font_size=20, bold=True, color=TEAL)
        y += Inches(0.4)
        txb = s.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.9))
        tf = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = body
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY
        y += Inches(1.1)

    # ── Slide 17: Assignment 5 ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📝 Assignment 5: To-Do List App (Final Project)")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Build a complete To-Do List App that:",
                font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate([
        "Lets the user add, view, mark complete, delete tasks",
        "Saves tasks to a text file (persists between runs!)",
        "Loads tasks from the file on startup",
        "Uses functions for each operation",
        "Handles errors gracefully",
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(19)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.75), Inches(11.5), Inches(0.45),
                "⭐ Bonus: Add due dates, priority levels, or categories.",
                font_size=18, italic=True, color=TEAL)
    add_textbox(s, Inches(0.8), Inches(5.25), Inches(11.5), Inches(0.4),
                "📁 Start with: starter-code/assignment5_starter.py",
                font_size=17, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(5.65), Inches(11.5), Inches(0.4),
                "📄 Full details: assignments/assignment5_todo_app.md",
                font_size=17, color=GRAY)

    # ── Slide 18: What's Next ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🚀 What's Next After This Course?")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.4),
                "Intermediate Topics to Explore:", font_size=20, bold=True, color=NAVY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(5.5), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "Classes & OOP — class, objects, inheritance",
        "List/dict comprehensions — advanced patterns",
        "Decorators & generators — powerful Python features",
        "Virtual environments — venv, pip",
        "Popular libraries: requests, pandas, flask",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {item}"
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(7.2), Inches(1.55), Inches(5.5), Inches(0.4),
                "Project Ideas:", font_size=20, bold=True, color=NAVY)

    txb2 = s.shapes.add_textbox(Inches(7.2), Inches(2.05), Inches(5.5), Inches(1.6))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for item in [
        "Web scraper with requests + BeautifulSoup",
        "Simple web app with Flask",
        "Data analysis with pandas",
        "Discord/Telegram bot",
    ]:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {item}"
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4),
                "Resources:", font_size=20, bold=True, color=NAVY)
    add_textbox(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
                "Real Python  realpython.com   |   Python Docs  docs.python.org/3/"
                "   |   Automate the Boring Stuff  automatetheboringstuff.com",
                font_size=16, color=TEAL)

    # ── Slide 19: Congratulations ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.0), Inches(11.33), Inches(1.2),
                "🎉 Congratulations!", font_size=56, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.3), Inches(11.33), Inches(0.7),
                "You've completed the Python for Beginners course!",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(2), Inches(3.1), Inches(9.33), Pt(2))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(3.3), Inches(11.33), Inches(0.4),
                "You now know how to:", font_size=20, bold=True,
                color=RGBColor(0xBF, 0xD7, 0xF5), align=PP_ALIGN.CENTER)

    achievements = [
        "Write and run Python programs",
        "Use variables, data types, and operators",
        "Control program flow with conditions and loops",
        "Write reusable functions",
        "Work with lists, tuples, and dictionaries",
        "Read and write files",
        "Handle errors gracefully",
        "Build a complete command-line application",
    ]

    txb = s.shapes.add_textbox(Inches(2.5), Inches(3.85), Inches(8.33), Inches(3.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in achievements:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"✅  {item}"
        run.font.size = Pt(17)
        run.font.color.rgb = WHITE

    add_textbox(s, Inches(1), Inches(7.0), Inches(11.33), Inches(0.4),
                "Keep coding, keep learning! 🐍",
                font_size=22, bold=True, color=RGBColor(0xA0, 0xC8, 0xFF),
                align=PP_ALIGN.CENTER)

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "session5_files_errors_project.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides → {out_path}")
