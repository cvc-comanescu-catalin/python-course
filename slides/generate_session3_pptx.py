"""Generate session3_functions.pptx from the Session 3 course content.

Usage:
    pip install python-pptx
    python slides/generate_session3_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)   # dark title backgrounds
TEAL   = RGBColor(0x00, 0x70, 0xC0)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # subtle slide background
GREEN  = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xC5, 0x56, 0x12)
GRAY   = RGBColor(0x44, 0x44, 0x44)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)  # VS-Code dark background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, bold=False, italic=False,
                color=GRAY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    """Full-width navy title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.4))
    fill_solid(bar, NAVY)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xF5)


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Dark VS-Code-style code block."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    fill_solid(box, CODE_BG)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)


def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # completely blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
                "🔧 Functions & Modular Thinking", font_size=48, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
                "Session 3: Writing Reusable Code",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(3), Inches(3.9), Inches(7.33), Pt(3))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                "Duration: 1 hour 40 minutes",
                font_size=20, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.5),
                "Write it once, use it everywhere! ♻️",
                font_size=18, italic=True, color=RGBColor(0xA0, 0xC8, 0xFF),
                align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📋 What We'll Cover Today")

    agenda = [
        "Recap & Assignment 2 Review",
        "Why Functions? The DRY Principle",
        "Defining Functions: def, parameters, return",
        "Default Parameters & Keyword Arguments",
        "Variable Scope: local vs. global",
        "Built-in Functions Cheat Sheet",
        "Importing Modules: math, random",
        "Practice",
    ]

    txb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(agenda, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}.  {item}"
        run.font.size = Pt(22)
        run.font.color.rgb = NAVY if i % 2 == 1 else TEAL

    # ── Slide 3: Recap & Assignment 2 Review ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🧠 Recap & Assignment 2 Review")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Quick check — what does each snippet output?",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# What does this output?\n"
                   "for i in range(3, 10, 3):\n"
                   "    print(i)\n"
                   "\n"
                   "# What does this do?\n"
                   "count = 0\n"
                   "while count < 5:\n"
                   "    if count == 3:\n"
                   "        break\n"
                   "    count += 1\n"
                   "print(count)\n"
                   "\n"
                   "# Fix the bug:\n"
                   'numbers = "12345"\n'
                   "for n in numbers:\n"
                   "    print(n * 2)    # What's the issue?",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.9),
                   font_size=14)

    # ── Slide 4: Why Functions? DRY Principle ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🤔 Why Functions? The DRY Principle")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "DRY = Don't Repeat Yourself",
                font_size=22, bold=True, color=NAVY)

    add_textbox(s, Inches(0.8), Inches(2.05), Inches(5.6), Inches(0.4),
                "❌ Without functions — repetitive:",
                font_size=18, bold=True, color=ORANGE)
    add_code_block(s,
                   'print("=" * 30)\n'
                   'print("  Welcome to the Game!")\n'
                   'print("=" * 30)\n'
                   "# ... lots of code ...\n"
                   'print("=" * 30)\n'
                   'print("  Game Over!")\n'
                   'print("=" * 30)',
                   Inches(0.8), Inches(2.5), Inches(5.6), Inches(2.8))

    add_textbox(s, Inches(7.0), Inches(2.05), Inches(5.6), Inches(0.4),
                "✅ With functions — clean:",
                font_size=18, bold=True, color=GREEN)
    add_code_block(s,
                   "def print_banner(message):\n"
                   '    print("=" * 30)\n'
                   '    print(f"  {message}")\n'
                   '    print("=" * 30)\n'
                   "\n"
                   'print_banner("Welcome to the Game!")\n'
                   "# ... lots of code ...\n"
                   'print_banner("Game Over!")',
                   Inches(7.0), Inches(2.5), Inches(5.6), Inches(2.8))

    add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                "Functions make code reusable, readable, and easier to maintain.",
                font_size=19, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ── Slide 5: Anatomy of a Function ───────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔧 Anatomy of a Function")

    add_code_block(s,
                   "#   keyword  name      parameters\n"
                   "#     \u2193       \u2193          \u2193\n"
                   'def greet(name, greeting="Hello"):\n'
                   '    """This is a docstring \u2014 explains what the function does."""\n'
                   '    message = f"{greeting}, {name}!"   # function body\n'
                   "    return message                      # return value\n"
                   "#      \u2191\n"
                   "#  return keyword\n"
                   "\n"
                   "# Calling the function:\n"
                   'result = greet("Alice")\n'
                   'print(result)           # Hello, Alice!\n'
                   "\n"
                   'result2 = greet("Bob", "Hi")\n'
                   "print(result2)          # Hi, Bob!",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=14)

    # ── Slide 6: Parameters vs Arguments ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📥 Parameters vs. Arguments")

    add_code_block(s,
                   "def add(a, b):       # a and b are PARAMETERS (in definition)\n"
                   "    return a + b\n"
                   "\n"
                   "result = add(3, 5)   # 3 and 5 are ARGUMENTS (in call)",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.8))

    add_textbox(s, Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.45),
                "Multiple parameters:", font_size=20, bold=True, color=NAVY)

    add_code_block(s,
                   "def describe_pet(name, animal_type, age):\n"
                   '    print(f"{name} is a {age}-year-old {animal_type}.")\n'
                   "\n"
                   'describe_pet("Rex", "dog", 3)\n'
                   "# Rex is a 3-year-old dog.\n"
                   "\n"
                   "# What if the function has no return?\n"
                   "def say_hi(name):\n"
                   '    print(f"Hi {name}!")\n'
                   "    # implicitly returns None\n"
                   "\n"
                   'value = say_hi("Alice")\n'
                   "print(value)   # None",
                   Inches(0.8), Inches(4.05), Inches(11.5), Inches(3.0))

    # ── Slide 7: Default Parameters & Keyword Arguments ───────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "⚙️ Default Parameters & Keyword Arguments")

    add_code_block(s,
                   "# Default parameters \u2014 used if argument is not provided\n"
                   "def power(base, exponent=2):\n"
                   "    return base ** exponent\n"
                   "\n"
                   "print(power(3))      # 9  \u2014 uses default exponent=2\n"
                   "print(power(3, 3))   # 27 \u2014 overrides default\n"
                   "\n"
                   "# Keyword arguments \u2014 order doesn't matter\n"
                   "def create_profile(name, age, city=\"Unknown\"):\n"
                   '    print(f"{name}, {age}, from {city}")\n'
                   "\n"
                   'create_profile("Alice", 25)\n'
                   'create_profile(age=30, name="Bob", city="Paris")\n'
                   'create_profile("Carol", city="Tokyo", age=22)',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=14)

    # ── Slide 8: Variable Scope ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔭 Variable Scope: Local vs. Global")

    add_code_block(s,
                   "total = 0   # Global variable \u2014 accessible everywhere\n"
                   "\n"
                   "def add_to_total(amount):\n"
                   "    global total          # Declare we're using the global\n"
                   "    total += amount       # Modify global variable\n"
                   "\n"
                   "def calculate():\n"
                   "    result = 100          # Local variable \u2014 only inside this function\n"
                   "    print(result)\n"
                   "\n"
                   "calculate()\n"
                   "# print(result)  \u2190 ERROR! result is local to calculate()\n"
                   "\n"
                   "add_to_total(50)\n"
                   "print(total)   # 50",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.8))

    add_textbox(s, Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.45),
                "Best practice: avoid global — return values instead!",
                font_size=19, bold=True, color=NAVY)

    add_code_block(s,
                   "def add(a, b):\n"
                   "    return a + b          # Return instead of using global\n"
                   "\n"
                   "total = add(10, 20)      # Assign the returned value",
                   Inches(0.8), Inches(6.05), Inches(11.5), Inches(1.2))

    # ── Slide 9: Built-in Functions Cheat Sheet ───────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📚 Built-in Functions Cheat Sheet")

    add_code_block(s,
                   "# Math\n"
                   "abs(-5)           # 5      \u2014 absolute value\n"
                   "round(3.7)        # 4      \u2014 round to nearest int\n"
                   "round(3.14159, 2) # 3.14   \u2014 round to 2 decimal places\n"
                   "max(1, 5, 3)      # 5      \u2014 maximum\n"
                   "min(1, 5, 3)      # 1      \u2014 minimum\n"
                   "sum([1, 2, 3])    # 6      \u2014 sum of list\n"
                   "\n"
                   "# Type & Info\n"
                   'type("hello")     # <class \'str\'>\n'
                   'len("Python")     # 6      \u2014 length\n'
                   "len([1, 2, 3])    # 3\n"
                   "\n"
                   "# Conversion\n"
                   'int("42")         # 42\n'
                   'float("3.14")     # 3.14\n'
                   "str(100)          # \"100\"\n"
                   "bool(0)           # False\n"
                   "list(\"abc\")       # ['a', 'b', 'c']",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=13)

    # ── Slide 10: Importing Modules ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📦 Importing Modules")

    add_code_block(s,
                   "# Import the whole module\n"
                   "import math\n"
                   "\n"
                   "print(math.pi)           # 3.141592653589793\n"
                   "print(math.sqrt(16))     # 4.0\n"
                   "print(math.floor(3.7))   # 3\n"
                   "print(math.ceil(3.2))    # 4\n"
                   "print(math.factorial(5)) # 120\n"
                   "\n"
                   "# Import specific things\n"
                   "from random import randint, choice, shuffle\n"
                   "\n"
                   "number = randint(1, 10)   # Random int between 1 and 10\n"
                   'item = choice(["a", "b", "c"])  # Random item from list',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=14)

    # ── Slide 11: The random Module ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎲 The random Module")

    add_code_block(s,
                   "import random\n"
                   "\n"
                   "# Random integer\n"
                   "random.randint(1, 6)       # Simulates a dice roll: 1-6\n"
                   "\n"
                   "# Random float between 0 and 1\n"
                   "random.random()            # e.g., 0.7234\n"
                   "\n"
                   "# Random choice from a sequence\n"
                   'fruits = ["apple", "banana", "cherry"]\n'
                   'random.choice(fruits)      # e.g., "banana"\n'
                   "\n"
                   "# Shuffle a list (modifies in place)\n"
                   "cards = [1, 2, 3, 4, 5]\n"
                   "random.shuffle(cards)      # e.g., [3, 1, 5, 2, 4]\n"
                   "\n"
                   "# Sample without replacement\n"
                   "random.sample(range(50), 6)  # Pick 6 lottery numbers",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=14)

    # ── Slide 12: Functions Calling Functions ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🧩 Functions Calling Functions")

    add_code_block(s,
                   "def square(n):\n"
                   "    return n ** 2\n"
                   "\n"
                   "def sum_of_squares(a, b):\n"
                   "    return square(a) + square(b)   # calls square() twice\n"
                   "\n"
                   "print(sum_of_squares(3, 4))  # 9 + 16 = 25\n"
                   "\n"
                   "# A complete example\n"
                   "def celsius_to_fahrenheit(c):\n"
                   "    return c * 9 / 5 + 32\n"
                   "\n"
                   "def temperature_report(city, temp_c):\n"
                   "    temp_f = celsius_to_fahrenheit(temp_c)\n"
                   '    print(f"{city}: {temp_c}\u00b0C = {temp_f:.1f}\u00b0F")\n'
                   "\n"
                   'temperature_report("Bucharest", 25)  # Bucharest: 25\u00b0C = 77.0\u00b0F\n'
                   'temperature_report("London", 15)     # London: 15\u00b0C = 59.0\u00b0F',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5),
                   font_size=14)

    # ── Slide 13: Practice Exercises ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎯 Practice Exercises")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Try these on your own:", font_size=22, bold=True, color=NAVY)

    exercises = [
        ("Exercise 1", "Write a function is_even(n) that returns True if n is even."),
        ("Exercise 2", "Write a function celsius_to_kelvin(c) that converts Celsius to Kelvin. (K = C + 273.15)"),
        ("Exercise 3", "Write a function count_vowels(text) that counts how many vowels are in a string."),
        ("Exercise 4", "Write a function roll_dice(sides=6) that simulates rolling a dice with a given number of sides (default: 6)."),
        ("Exercise 5", "Write a function bmi(weight_kg, height_m) that calculates Body Mass Index: weight / height\u00b2."),
    ]

    y = Inches(2.2)
    for title, body in exercises:
        add_textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.4),
                    title, font_size=19, bold=True, color=TEAL)
        y += Inches(0.4)
        txb = s.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.55))
        tf = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = body
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY
        y += Inches(0.65)

    # ── Slide 14: Assignment 3 ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📝 Assignment 3: Calculator with Functions")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Due: Before Session 4",
                font_size=20, bold=True, color=ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.45),
                "Build a command-line calculator that:", font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(11.5), Inches(2.2))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate([
        "Has separate functions: add(), subtract(), multiply(), divide()",
        "Shows a menu to the user and asks them to pick an operation",
        "Handles division by zero gracefully",
        'Loops until the user types "quit"',
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(18)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.85), Inches(11.5), Inches(0.45),
                "⭐ Bonus:",
                font_size=19, bold=True, color=TEAL)

    txb2 = s.shapes.add_textbox(Inches(0.8), Inches(5.35), Inches(11.5), Inches(0.9))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in [
        "Add power() and square_root() functions using the math module",
        "Add a history of calculations",
    ]:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {b}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
                "📁 Start with: starter-code/assignment3_starter.py",
                font_size=17, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
                "📄 Full details: assignments/assignment3_calculator.md",
                font_size=17, color=GRAY)

    # ── Slide 15: Session 3 Recap ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎉 Session 3 Recap")

    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "Why functions make code better (DRY principle)",
        "Defining functions with  def",
        "Parameters, arguments, and  return",
        "Default parameters & keyword arguments",
        "Local vs. global scope",
        "Useful built-in functions",
        "Importing and using  math  and  random",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  \u2705  {item}"
        run.font.size = Pt(21)
        run.font.color.rgb = NAVY

    add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "See you in Session 4! 🚀",
                font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "session3_functions.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides → {out_path}")
