"""Generate session4_data_structures.pptx from the Session 4 course content.

Usage:
    pip install python-pptx
    python slides/generate_session4_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)   # dark title backgrounds
TEAL   = RGBColor(0x00, 0x70, 0xC0)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # subtle slide background
GREEN  = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xC5, 0x56, 0x12)
GRAY   = RGBColor(0x44, 0x44, 0x44)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)  # VS-Code dark background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, bold=False, italic=False,
                color=GRAY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    """Full-width navy title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.4))
    fill_solid(bar, NAVY)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xF5)


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Dark VS-Code-style code block."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    fill_solid(box, CODE_BG)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)


def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # completely blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
                "📦 Data Structures", font_size=52, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
                "Session 4: Lists, Tuples & Dictionaries",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(3), Inches(3.9), Inches(7.33), Pt(3))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                "Duration: 1 hour 40 minutes",
                font_size=20, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📋 What We'll Cover Today")

    agenda = [
        "Recap & Assignment 3 Review",
        "Lists: creation, indexing, slicing, methods",
        "Looping through lists",
        "List comprehensions (intro)",
        "Tuples — immutable sequences",
        "Dictionaries: key-value pairs, methods",
        "Nested data structures",
        "Practice",
    ]

    txb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(agenda, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}.  {item}"
        run.font.size = Pt(22)
        run.font.color.rgb = NAVY if i % 2 == 1 else TEAL

    # ── Slide 3: Recap & Assignment 3 Review ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🧠 Recap & Assignment 3 Review")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Quick check:", font_size=20, bold=True, color=NAVY)

    add_code_block(s,
                   "# What does this print?\n"
                   "def mystery(x, y=10):\n"
                   "    return x * y\n"
                   "\n"
                   "print(mystery(3))\n"
                   "print(mystery(3, 5))\n"
                   "\n"
                   "# What's the scope issue here?\n"
                   "def double():\n"
                   "    result = x * 2   # Is this okay?\n"
                   "\n"
                   "x = 5\n"
                   "double()",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.6))

    # ── Slide 4: Lists — Ordered Collections ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📋 Lists — Ordered Collections")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "A list is an ordered, changeable collection of items.",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# Creating lists\n"
                   'fruits = ["apple", "banana", "cherry"]\n'
                   "numbers = [1, 2, 3, 4, 5]\n"
                   'mixed = [42, "hello", True, 3.14]\n'
                   "empty = []\n"
                   "\n"
                   "# Indexing — starts at 0!\n"
                   "#   Index:  0        1         2\n"
                   'fruits = ["apple", "banana", "cherry"]\n'
                   "\n"
                   'print(fruits[0])    # "apple"\n'
                   'print(fruits[1])    # "banana"\n'
                   'print(fruits[-1])   # "cherry"  <- negative index from end\n'
                   'print(fruits[-2])   # "banana"',
                   Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.7))

    # ── Slide 5: List Slicing ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "✂️ List Slicing")

    add_code_block(s,
                   "numbers = [0, 1, 2, 3, 4, 5, 6, 7, 8, 9]\n"
                   "\n"
                   "# [start:stop]  — stop is NOT included\n"
                   "print(numbers[2:5])    # [2, 3, 4]\n"
                   "print(numbers[:4])     # [0, 1, 2, 3]  — from beginning\n"
                   "print(numbers[6:])     # [6, 7, 8, 9]  — to end\n"
                   "print(numbers[:])      # [0, 1, 2, ... 9]  — full copy\n"
                   "\n"
                   "# [start:stop:step]\n"
                   "print(numbers[::2])    # [0, 2, 4, 6, 8]  — every 2nd\n"
                   "print(numbers[::-1])   # [9, 8, 7, ... 0]  — reversed!\n"
                   "\n"
                   "# Visual:\n"
                   "# [0, 1, 2, 3, 4, 5, 6, 7, 8, 9]\n"
                   "#  ^\n"
                   "# [2:5] -> takes index 2, 3, 4 -> [2, 3, 4]",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 6: Common List Methods ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🛠️ Common List Methods")

    add_code_block(s,
                   'fruits = ["apple", "banana"]\n'
                   "\n"
                   "# Adding items\n"
                   'fruits.append("cherry")        # Add to end -> [..., "cherry"]\n'
                   'fruits.insert(1, "blueberry")  # Insert at index 1\n'
                   "\n"
                   "# Removing items\n"
                   'fruits.remove("banana")        # Remove by value\n'
                   "last = fruits.pop()            # Remove & return last item\n"
                   "item = fruits.pop(0)           # Remove & return at index 0\n"
                   "\n"
                   "# Other useful methods\n"
                   "numbers = [3, 1, 4, 1, 5, 9, 2, 6]\n"
                   "numbers.sort()                 # Sort in place -> [1, 1, 2, 3, 4, 5, 6, 9]\n"
                   "numbers.reverse()              # Reverse in place\n"
                   "print(len(numbers))            # 8\n"
                   "print(numbers.count(1))        # 2 — how many times 1 appears\n"
                   "print(numbers.index(5))        # 4 — index of value 5",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 7: Looping Through Lists ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔁 Looping Through Lists")

    add_code_block(s,
                   'fruits = ["apple", "banana", "cherry", "date"]\n'
                   "\n"
                   "# Simple loop\n"
                   "for fruit in fruits:\n"
                   "    print(fruit)\n"
                   "\n"
                   "# Loop with index using enumerate()\n"
                   "for i, fruit in enumerate(fruits):\n"
                   '    print(f"{i}: {fruit}")\n'
                   "# 0: apple\n"
                   "# 1: banana\n"
                   "# 2: cherry\n"
                   "\n"
                   "# Loop and modify (use range + len)\n"
                   "numbers = [1, 2, 3, 4, 5]\n"
                   "for i in range(len(numbers)):\n"
                   "    numbers[i] = numbers[i] * 2\n"
                   "print(numbers)   # [2, 4, 6, 8, 10]",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 8: List Comprehensions ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "✨ List Comprehensions")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "A concise way to create or transform lists:",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# Traditional way:\n"
                   "squares = []\n"
                   "for n in range(1, 6):\n"
                   "    squares.append(n ** 2)\n"
                   "\n"
                   "# List comprehension — same result in one line!\n"
                   "squares = [n ** 2 for n in range(1, 6)]\n"
                   "print(squares)   # [1, 4, 9, 16, 25]\n"
                   "\n"
                   "# With a condition (filter):\n"
                   "evens = [n for n in range(20) if n % 2 == 0]\n"
                   "print(evens)     # [0, 2, 4, 6, 8, 10, 12, 14, 16, 18]\n"
                   "\n"
                   "# Transform strings:\n"
                   'words = ["hello", "world", "python"]\n'
                   "upper = [word.upper() for word in words]\n"
                   "print(upper)     # ['HELLO', 'WORLD', 'PYTHON']",
                   Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.7))

    # ── Slide 9: Tuples — Immutable Sequences ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔒 Tuples — Immutable Sequences")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Tuples are like lists, but cannot be changed after creation.",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# Creating tuples (use parentheses)\n"
                   "point = (3, 7)\n"
                   "rgb = (255, 128, 0)\n"
                   "single = (42,)   # <- Note the comma! (42) is just an int\n"
                   "\n"
                   "# Accessing elements — same as lists\n"
                   "print(point[0])   # 3\n"
                   "print(rgb[-1])    # 0\n"
                   "\n"
                   "# Unpacking\n"
                   "x, y = point\n"
                   'print(f"x={x}, y={y}")   # x=3, y=7\n'
                   "\n"
                   "lat, lon = 44.4268, 26.1025   # Bucharest coordinates",
                   Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.6))

    add_textbox(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.45),
                "When to use tuples vs. lists:", font_size=19, bold=True, color=NAVY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.9))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in ["Tuple: Fixed data (coordinates, RGB, days of week)",
                 "List: Data that will change (shopping cart, user list)"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {item}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    # ── Slide 10: Dictionaries — Key-Value Pairs ──────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📖 Dictionaries — Key-Value Pairs")

    add_code_block(s,
                   "# Creating a dictionary\n"
                   "student = {\n"
                   '    "name": "Alice",\n'
                   '    "age": 20,\n'
                   '    "grade": "A",\n'
                   '    "courses": ["Math", "Physics"]\n'
                   "}\n"
                   "\n"
                   "# Accessing values\n"
                   'print(student["name"])      # "Alice"\n'
                   "print(student.get(\"age\"))   # 20\n"
                   'print(student.get("gpa", 0.0))  # 0.0 (default if key missing)\n'
                   "\n"
                   "# Adding / updating\n"
                   'student["email"] = "alice@example.com"  # Add new key\n'
                   'student["age"] = 21                     # Update existing key\n'
                   "\n"
                   "# Deleting\n"
                   'del student["grade"]\n'
                   'removed = student.pop("age")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 11: Dictionary Methods ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔍 Dictionary Methods")

    add_code_block(s,
                   'person = {"name": "Bob", "age": 30, "city": "Paris"}\n'
                   "\n"
                   "# Keys, values, items\n"
                   "print(person.keys())    # dict_keys(['name', 'age', 'city'])\n"
                   "print(person.values())  # dict_values(['Bob', 30, 'Paris'])\n"
                   "print(person.items())   # dict_items([('name', 'Bob'), ...])\n"
                   "\n"
                   "# Looping\n"
                   "for key in person:\n"
                   '    print(key, "->", person[key])\n'
                   "\n"
                   "for key, value in person.items():\n"
                   '    print(f"{key}: {value}")\n'
                   "\n"
                   "# Check if key exists\n"
                   'if "name" in person:\n'
                   '    print("Name is:", person["name"])',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 12: Nested Data Structures ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🏗️ Nested Data Structures")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Combine lists and dicts for complex data:",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# List of dictionaries (common pattern!)\n"
                   "students = [\n"
                   '    {"name": "Alice", "grade": 92},\n'
                   '    {"name": "Bob",   "grade": 85},\n'
                   '    {"name": "Carol", "grade": 97},\n'
                   "]\n"
                   "\n"
                   "# Access nested data\n"
                   'print(students[0]["name"])   # "Alice"\n'
                   "print(students[2][\"grade\"])  # 97\n"
                   "\n"
                   "# Loop through list of dicts\n"
                   "for student in students:\n"
                   "    print(f\"{student['name']}: {student['grade']}\")\n"
                   "\n"
                   "# Dictionary of lists\n"
                   "gradebook = {\n"
                   '    "Alice": [90, 85, 92],\n'
                   '    "Bob":   [78, 82, 80],\n'
                   "}\n"
                   'print(sum(gradebook["Alice"]) / len(gradebook["Alice"]))  # Average',
                   Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.9))

    # ── Slide 13: Practice Exercises ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎯 Practice Exercises")

    exercises = [
        ("Exercise 1",
         "Create a list of 5 of your favorite movies. Print them numbered (1. Movie Name)."),
        ("Exercise 2",
         "Use a list comprehension to create a list of all numbers divisible by 3 between 1 and 50."),
        ("Exercise 3",
         "Create a dictionary for a contact (name, phone, email, city). Print each field on its own line."),
        ("Exercise 4",
         "Given scores = [78, 92, 85, 90, 88, 76], calculate the average without using sum() (use a loop)."),
        ("Exercise 5",
         "Write a function most_common(lst) that returns the most frequent item in a list."),
    ]

    y = Inches(1.6)
    for title, body in exercises:
        add_textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.4),
                    title, font_size=19, bold=True, color=TEAL)
        y += Inches(0.38)
        txb = s.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.5))
        tf = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = body
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY
        y += Inches(0.65)

    # ── Slide 14: Assignment 4 ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📝 Assignment 4: Student Gradebook")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Due: Before Session 5",
                font_size=20, bold=True, color=ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.45),
                "Create a program that:", font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.1))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate([
        "Stores student names and their grades in a dictionary",
        "Lets the user: add a student, view all students, look up a student's grade, calculate the class average",
        "Uses a menu-driven loop",
        "Validates user input (what if the student doesn't exist?)",
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.45),
                "⭐ Bonus:", font_size=18, bold=True, color=TEAL)

    txb2 = s.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.2))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in [
        "Store multiple grades per student (dict of lists)",
        "Find highest/lowest scoring student",
        "Show a grade distribution (A/B/C/D/F count)",
        "Allow removing a student",
    ]:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {b}"
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.4),
                "📁 Start with: starter-code/assignment4_starter.py",
                font_size=17, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.4),
                "📄 Full details: assignments/assignment4_gradebook.md",
                font_size=17, color=GRAY)

    # ── Slide 15: Recap ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎉 Session 4 Recap")

    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "Lists: creation, indexing, slicing, methods",
        "Looping through lists with for and enumerate()",
        "List comprehensions for concise list creation",
        "Tuples: immutable, use for fixed data",
        "Dictionaries: key-value pairs, CRUD operations",
        "Iterating dicts with .keys(), .values(), .items()",
        "Nested structures: list of dicts, dict of lists",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  \u2705  {item}"
        run.font.size = Pt(21)
        run.font.color.rgb = NAVY

    add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "See you in Session 5 — the final session! 🚀",
                font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "session4_data_structures.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides -> {out_path}")
