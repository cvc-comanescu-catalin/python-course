"""Generate session1_intro_and_basics.pptx from the Session 1 course content.

Usage:
    pip install python-pptx
    python slides/generate_session1_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)   # dark title backgrounds
TEAL   = RGBColor(0x00, 0x70, 0xC0)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # subtle slide background
GREEN  = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xC5, 0x56, 0x12)
GRAY   = RGBColor(0x44, 0x44, 0x44)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)  # VS-Code dark background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, bold=False, italic=False,
                color=GRAY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    """Full-width navy title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.4))
    fill_solid(bar, NAVY)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xF5)


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Dark VS-Code-style code block."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    fill_solid(box, CODE_BG)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)


def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # completely blank layout

    # ── Slide 1: Title / Welcome ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
                "🐍 Welcome to Python!", font_size=52, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
                "Session 1: Introduction to Python & Programming Basics",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(3), Inches(3.9), Inches(7.33), Pt(3))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                "Duration: 1 hour 40 minutes",
                font_size=20, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.5),
                "Your instructor is excited to meet you! 🎉",
                font_size=18, italic=True, color=RGBColor(0xA0, 0xC8, 0xFF),
                align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📋 What We'll Cover Today")

    agenda = [
        "What is Python & why learn it?",
        "Setting up your environment",
        "Hello, World!",
        "Variables & Data Types",
        "Arithmetic & Assignment Operators",
        "User Input with input()",
        "Type Conversion",
        "Comments & Code Readability",
        "Live Coding & Practice",
    ]

    txb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(agenda, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}.  {item}"
        run.font.size = Pt(22)
        run.font.color.rgb = NAVY if i % 2 == 1 else TEAL

    # ── Slide 3: What is Python? ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🌍 What is Python?")

    bullets = [
        "A general-purpose, high-level programming language",
        "Created by Guido van Rossum in 1991",
        "Named after Monty Python's Flying Circus 🎭",
        "One of the most popular languages in the world",
    ]
    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.5), Inches(3))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(20)
        run.font.color.rgb = GRAY

    add_code_block(s,
                   "Timeline:\n"
                   "1991 ──► Python 1.0\n"
                   "1994 ──► Python 1.4\n"
                   "2000 ──► Python 2.0\n"
                   "2008 ──► Python 3.0\n"
                   "2023 ──► Python 3.12 (current)",
                   Inches(7.5), Inches(1.6), Inches(5.2), Inches(3.2))

    # ── Slide 4: What Can Python Do? ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🚀 What Can Python Do?")

    table_data = [
        ("🌐 Web Development",       "Instagram, Pinterest, Django"),
        ("🤖 AI & Machine Learning", "TensorFlow, ChatGPT tools"),
        ("📊 Data Science",           "Pandas, NumPy, Jupyter"),
        ("🎮 Game Development",       "Pygame"),
        ("🔬 Science & Research",     "NASA, CERN"),
        ("🤖 Automation",             "Scripts, bots, task automation"),
    ]

    tbl = s.shapes.add_table(
        len(table_data) + 1, 2,
        Inches(1), Inches(1.6), Inches(11.33), Inches(4.2),
    ).table
    tbl.columns[0].width = Inches(4.8)
    tbl.columns[1].width = Inches(6.53)

    for ci, hdr in enumerate(["Field", "Examples"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE

    for ri, (field, example) in enumerate(table_data, 1):
        bg = LIGHT if ri % 2 == 1 else WHITE
        for ci, txt in enumerate([field, example]):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = txt
            run.font.size = Pt(16)
            run.font.color.rgb = GRAY

    add_textbox(s, Inches(1), Inches(6.1), Inches(11.33), Inches(0.5),
                "Python is everywhere — and it's beginner-friendly!",
                font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ── Slide 5: Setting Up ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "💻 Setting Up Your Environment")

    add_textbox(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.5),
                "Option 1: Local Install (Recommended)",
                font_size=20, bold=True, color=NAVY)

    txb = s.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.8), Inches(1.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate(["Go to python.org/downloads",
                            "Download Python 3.10 or newer",
                            "Install VS Code + Python extension"], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.6), Inches(3.7), Inches(5.8), Inches(0.5),
                "Option 2: Online (Zero Setup)",
                font_size=20, bold=True, color=NAVY)

    txb2 = s.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(5.8), Inches(1.2))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for t in ["Google Colab — free, in your browser", "repl.it — online IDE"]:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {t}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.5),
                "Verify your installation:",
                font_size=18, bold=True, color=NAVY)
    add_code_block(s, "python --version\n# Should show: Python 3.10.x or higher",
                   Inches(7.2), Inches(2.1), Inches(5.5), Inches(1.3))

    # ── Slide 6: Hello, World! ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "👋 Your First Python Program")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Open your editor and type:", font_size=20, color=GRAY)
    add_code_block(s, 'print("Hello, World!")',
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.9))
    add_textbox(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.5),
                "Run it! You should see:", font_size=20, color=GRAY)
    add_code_block(s, "Hello, World!",
                   Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.7))
    add_textbox(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
                "Congratulations — you're a programmer! 🎉",
                font_size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                "More print examples:", font_size=20, bold=True, color=NAVY)
    add_code_block(s,
                   'print("My name is Alice")\n'
                   'print("2 + 2 =", 2 + 2)\n'
                   'print("Python is", "awesome!")',
                   Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.4))

    # ── Slide 7: Variables ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "📦 Variables — Boxes That Hold Values")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "A variable is a named storage location for data.",
                font_size=20, color=GRAY)
    add_code_block(s,
                   'name       = "Alice"   # str   — text\n'
                   "age        = 25        # int   — whole number\n"
                   "height     = 1.68      # float — decimal number\n"
                   "is_student = True      # bool  — True or False",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.8))

    add_textbox(s, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.5),
                "Rules for variable names:", font_size=20, bold=True, color=NAVY)
    txb = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for r in ["Start with a letter or _ (not a number)",
              "No spaces — use _ instead:  first_name",
              "Case-sensitive:  age  ≠  Age"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {r}"
        run.font.size = Pt(19)
        run.font.color.rgb = GRAY

    # ── Slide 8: Data Types ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔠 Data Types")

    add_code_block(s,
                   "# Integer — whole numbers\n"
                   "score = 100\n"
                   "year  = 2024\n\n"
                   "# Float — decimal numbers\n"
                   "price = 9.99\n"
                   "pi    = 3.14159\n\n"
                   "# String — text (use quotes)\n"
                   'greeting = "Hello!"\n'
                   "city     = 'Bucharest'\n\n"
                   "# Boolean — True or False\n"
                   "is_raining  = False\n"
                   "passed_exam = True\n\n"
                   "# Check a type with type()\n"
                   "print(type(score))       # <class 'int'>\n"
                   "print(type(price))       # <class 'float'>\n"
                   "print(type(greeting))    # <class 'str'>\n"
                   "print(type(is_raining))  # <class 'bool'>",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4),
                   font_size=14)

    # ── Slide 9: Arithmetic Operators ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "➕ Arithmetic Operators")

    add_code_block(s,
                   "a = 10\n"
                   "b = 3\n\n"
                   "print(a + b)   # 13     Addition\n"
                   "print(a - b)   # 7      Subtraction\n"
                   "print(a * b)   # 30     Multiplication\n"
                   "print(a / b)   # 3.333  Division (always float)\n"
                   "print(a // b)  # 3      Floor Division\n"
                   "print(a % b)   # 1      Modulo (remainder)\n"
                   "print(a ** b)  # 1000   Exponentiation (10³)",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.6))

    add_textbox(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
                "Order of Operations: PEMDAS / BODMAS",
                font_size=20, bold=True, color=NAVY)
    add_code_block(s,
                   "result = 2 + 3 * 4    # 14  (not 20!)\n"
                   "result = (2 + 3) * 4  # 20",
                   Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2))

    # ── Slide 10: Assignment Operators ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔄 Assignment Operators")

    add_code_block(s,
                   "x = 10     # Assign\n"
                   "x += 5     # x = x + 5   → 15\n"
                   "x -= 3     # x = x - 3   → 12\n"
                   "x *= 2     # x = x * 2   → 24\n"
                   "x /= 4     # x = x / 4   → 6.0\n"
                   "x //= 2    # x = x // 2  → 3.0\n"
                   "x **= 3    # x = x ** 3  → 27.0",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.2))

    add_textbox(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
                "String operations:", font_size=20, bold=True, color=NAVY)
    add_code_block(s,
                   'first    = "Hello"\n'
                   'last     = "World"\n'
                   'full     = first + " " + last   # "Hello World"  — concatenation\n'
                   'repeated = "Ha" * 3              # "HaHaHa"       — repetition',
                   Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.6))

    # ── Slide 11: User Input ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "⌨️ User Input with input()")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "input()  pauses the program and waits for the user to type something.",
                font_size=20, color=GRAY)
    add_code_block(s,
                   'name = input("What is your name? ")\n'
                   'print("Hello,", name)',
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.2))
    add_textbox(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
                "⚠  Important:  input()  always returns a string!",
                font_size=20, bold=True, color=ORANGE)
    add_code_block(s,
                   'age_str = input("How old are you? ")\n'
                   "print(type(age_str))   # <class 'str'>  ← NOT int!\n\n"
                   "# You need to convert it:\n"
                   "age = int(age_str)\n"
                   "print(type(age))       # <class 'int'>  ✓",
                   Inches(0.8), Inches(3.9), Inches(11.5), Inches(2.8))

    # ── Slide 12: Type Conversion ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🔁 Type Conversion")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Convert between types using built-in functions:",
                font_size=20, color=GRAY)
    add_code_block(s,
                   "# String → number\n"
                   'age   = int("25")         # 25\n'
                   'price = float("9.99")     # 9.99\n\n'
                   "# Number → string\n"
                   'text = str(42)             # "42"\n'
                   'text = str(3.14)           # "3.14"\n\n'
                   "# One-liner (common pattern):\n"
                   'age   = int(input("Enter your age: "))\n'
                   'price = float(input("Enter price: "))\n\n'
                   "# Boolean conversions\n"
                   "bool(0)    # False\n"
                   "bool(1)    # True\n"
                   'bool("")   # False\n'
                   'bool("hi") # True',
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.9),
                   font_size=14)

    # ── Slide 13: Comments ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "💬 Comments & Code Readability")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Comments explain your code — they are ignored by Python.",
                font_size=20, color=GRAY)
    add_code_block(s,
                   "# This is a single-line comment\n\n"
                   "# Calculate the area of a rectangle\n"
                   "width  = 5\n"
                   "height = 3\n"
                   "area   = width * height   # multiply width by height\n"
                   "print(area)               # prints 15\n\n"
                   "# Multi-line comment (use multiple #):\n"
                   "# Step 1: Get user input\n"
                   "# Step 2: Do calculation\n"
                   "# Step 3: Display result",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(3.4))

    add_textbox(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
                "Best practices:", font_size=20, bold=True, color=NAVY)
    txb = s.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in ["Write comments for why, not just what",
              "Use descriptive variable names:  user_age  not  x",
              "Keep lines under ~79 characters"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  • {b}"
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY

    # ── Slide 14: Putting It All Together ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🧮 Putting It All Together")

    add_code_block(s,
                   '# A simple "About Me" program\n\n'
                   "# Get user information\n"
                   'name  = input("Enter your name: ")\n'
                   'age   = int(input("Enter your age: "))\n'
                   'hobby = input("Enter your favorite hobby: ")\n\n'
                   "# Calculate birth year\n"
                   "current_year = 2024\n"
                   "birth_year   = current_year - age\n\n"
                   "# Display formatted output\n"
                   "print()\n"
                   'print("=== About Me ===")\n'
                   "print(f\"Hi, I'm {name}!\")\n"
                   "print(f\"I'm {age} years old, born around {birth_year}.\")\n"
                   "print(f\"My favorite hobby is {hobby}.\")",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.6),
                   font_size=15)

    # ── Slide 15: Practice Exercises ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎯 Practice Exercises")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Try these on your own:", font_size=22, bold=True, color=NAVY)

    exercises = [
        ("Exercise 1", "Print your name and age using variables"),
        ("Exercise 2",
         "Calculate and print the area of a circle\n"
         "  • Get the radius from the user\n"
         "  • Use  pi = 3.14159\n"
         "  • Area = π × r²"),
        ("Exercise 3",
         "Temperature converter\n"
         "  • Ask for temperature in Celsius\n"
         "  • Convert to Fahrenheit:  F = C × 9/5 + 32\n"
         "  • Print the result"),
    ]

    y = Inches(2.2)
    for title, body in exercises:
        add_textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.4),
                    title, font_size=20, bold=True, color=TEAL)
        y += Inches(0.4)
        txb = s.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.9))
        tf = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = body
        run.font.size = Pt(17)
        run.font.color.rgb = GRAY
        y += Inches(1.05)

    # ── Slide 16: Assignment ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, '📝 Assignment 1: "About Me" Script')

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Due: Before Session 2",
                font_size=20, bold=True, color=ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.45),
                "Write a Python script that:", font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate([
        "Asks the user for their name, age, and favorite hobby",
        "Calculates the year they were born (approximately)",
        "Prints a formatted summary:\n"
        '   "Hi, I\'m Alice. I\'m 25 years old, born around 2001, and I love hiking!"',
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(18)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.45),
                "⭐ Bonus: Ask for two numbers and print their sum, difference, product, and quotient.",
                font_size=18, italic=True, color=TEAL)
    add_textbox(s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.4),
                "📁 Start with:  starter-code/assignment1_starter.py",
                font_size=17, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
                "📄 Full details:  assignments/assignment1_about_me.md",
                font_size=17, color=GRAY)

    # ── Slide 17: Recap ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "🎉 Session 1 Recap")

    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "What Python is and why it's awesome",
        "How to set up your coding environment",
        "print()  — displaying output",
        "Variables and the 4 basic data types",
        "Arithmetic and assignment operators",
        "input()  — getting user input",
        "Type conversion:  int(),  float(),  str()",
        "Writing comments",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  ✅  {item}"
        run.font.size = Pt(21)
        run.font.color.rgb = NAVY

    add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "See you in Session 2! 🚀",
                font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "session1_intro_and_basics.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides → {out_path}")
