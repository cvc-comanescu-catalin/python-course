"""Generate session2_control_flow.pptx from the Session 2 course content.

Usage:
    pip install python-pptx
    python slides/generate_session2_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)   # dark title backgrounds
TEAL   = RGBColor(0x00, 0x70, 0xC0)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # subtle slide background
GREEN  = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xC5, 0x56, 0x12)
GRAY   = RGBColor(0x44, 0x44, 0x44)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)  # VS-Code dark background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_textbox(slide, left, top, width, height,
                text="", font_size=18, bold=False, italic=False,
                color=GRAY, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_title_bar(slide, title, subtitle=None):
    """Full-width navy title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.4))
    fill_solid(bar, NAVY)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xF5)


def add_code_block(slide, code, left, top, width, height, font_size=13):
    """Dark VS-Code-style code block."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    fill_solid(box, CODE_BG)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.splitlines():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)


def set_slide_bg(slide, colour):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # completely blank layout

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)

    hero = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(hero, NAVY)
    hero.line.fill.background()

    add_textbox(s, Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
                "\U0001f500 Control Flow", font_size=52, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
                "Session 2: Making Decisions & Repeating Actions",
                font_size=26, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    line_box = s.shapes.add_shape(1, Inches(3), Inches(3.9), Inches(7.33), Pt(3))
    fill_solid(line_box, TEAL)
    line_box.line.fill.background()

    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                "Duration: 1 hour 40 minutes",
                font_size=20, color=RGBColor(0xBF, 0xD7, 0xF5),
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.5),
                "Get ready to make your programs think and repeat! \U0001f389",
                font_size=18, italic=True, color=RGBColor(0xA0, 0xC8, 0xFF),
                align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f4cb What We'll Cover Today")

    agenda = [
        "Quick Recap Quiz",
        "Comparison Operators",
        "if, elif, else statements",
        "Logical Operators: and, or, not",
        "while loops",
        "for loops & range()",
        "break and continue",
        "Live Coding: FizzBuzz & Guessing Game",
    ]

    txb = s.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(agenda, 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}.  {item}"
        run.font.size = Pt(22)
        run.font.color.rgb = NAVY if i % 2 == 1 else TEAL

    # ── Slide 3: Quick Recap Quiz ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f9e0 Quick Recap Quiz")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "Answer these in your head (or type them out!):",
                font_size=20, color=GRAY)

    questions = [
        '1. What does  type("hello")  return?',
        "2. What is the result of  10 % 3 ?",
        "3. What does  input()  always return?",
        "4. Fix this code:",
    ]

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for q in questions:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = q
        run.font.size = Pt(20)
        run.font.color.rgb = GRAY

    add_code_block(s,
                   'age = input("Your age: ")\n'
                   "next_year = age + 1   # \u2190 This crashes! Why?",
                   Inches(1.5), Inches(4.3), Inches(10), Inches(1.3))

    # ── Slide 4: Comparison Operators ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\u2696\ufe0f Comparison Operators")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "These return  True  or  False :",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "x = 10\n"
                   "y = 5\n\n"
                   "print(x == y)   # False  \u2014 Equal to\n"
                   "print(x != y)   # True   \u2014 Not equal to\n"
                   "print(x > y)    # True   \u2014 Greater than\n"
                   "print(x < y)    # False  \u2014 Less than\n"
                   "print(x >= 10)  # True   \u2014 Greater than or equal\n"
                   "print(x <= 9)   # False  \u2014 Less than or equal",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.9))

    add_code_block(s,
                   "# Works with strings too!\n"
                   'name = "Alice"\n'
                   'print(name == "Alice")   # True\n'
                   'print(name == "alice")   # False  \u2190 case-sensitive!',
                   Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.8))

    # ── Slide 5: if / elif / else ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f500 if / elif / else")

    # ASCII flowchart in monospace textbox
    flowchart = (
        "              \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "              \u2502  condition? \u2502\n"
        "              \u2514\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "                     \u2502\n"
        "         True \u25c4\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u25ba False\n"
        "         \u2502                   \u2502\n"
        "   \u250c\u2500\u2500\u2500\u2500\u2500\u25bc\u2500\u2500\u2500\u2500\u2510     \u250c\u2500\u2500\u2500\u2500\u2500\u25bc\u2500\u2500\u2500\u2500\u2510\n"
        "   \u2502 do this    \u2502     \u2502  do that   \u2502\n"
        "   \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518     \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518"
    )

    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(5.8), Inches(2.6))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in flowchart.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.name = "Courier New"
        run.font.color.rgb = NAVY

    add_code_block(s,
                   'age = int(input("Enter your age: "))\n\n'
                   "if age < 13:\n"
                   '    print("You are a child.")\n'
                   "elif age < 18:\n"
                   '    print("You are a teenager.")\n'
                   "elif age < 65:\n"
                   '    print("You are an adult.")\n'
                   "else:\n"
                   '    print("You are a senior.")',
                   Inches(7.0), Inches(1.55), Inches(5.7), Inches(3.5))

    # ── Slide 6: More if / elif / else Examples ───────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f500 More if / elif / else Examples")

    add_code_block(s,
                   "# Grade classifier\n"
                   'score = int(input("Enter your score: "))\n\n'
                   "if score >= 90:\n"
                   '    grade = "A"\n'
                   "elif score >= 80:\n"
                   '    grade = "B"\n'
                   "elif score >= 70:\n"
                   '    grade = "C"\n'
                   "elif score >= 60:\n"
                   '    grade = "D"\n'
                   "else:\n"
                   '    grade = "F"\n\n'
                   'print(f"Your grade is: {grade}")',
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8))

    add_textbox(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5),
                "\u2605 Key rule: Conditions are checked top to bottom \u2014 first True wins!",
                font_size=20, bold=True, color=NAVY)

    # ── Slide 7: Logical Operators ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f9e9 Logical Operators")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Combine multiple conditions:",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# and \u2014 both must be True\n"
                   "age = 20\n"
                   "has_id = True\n"
                   "print(age >= 18 and has_id)   # True\n\n"
                   "# or \u2014 at least one must be True\n"
                   "is_weekend = False\n"
                   "is_holiday = True\n"
                   "print(is_weekend or is_holiday)   # True\n\n"
                   "# not \u2014 reverses True/False\n"
                   "is_raining = False\n"
                   "print(not is_raining)   # True",
                   Inches(0.8), Inches(2.1), Inches(7.0), Inches(3.8))

    # Truth table for "and" as a textbox
    add_textbox(s, Inches(8.2), Inches(2.0), Inches(4.8), Inches(0.45),
                "Truth Table for  and :", font_size=18, bold=True, color=NAVY)

    truth_and = (
        "  A       B     A and B\n"
        "  True    True    True\n"
        "  True    False   False\n"
        "  False   True    False\n"
        "  False   False   False"
    )
    txb = s.shapes.add_textbox(Inches(8.2), Inches(2.5), Inches(4.8), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in truth_and.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = "Courier New"
        run.font.color.rgb = GRAY

    # ── Slide 8: Truth Table for or and not ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f504 Truth Table for `or` and `not`")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(5.5), Inches(0.4),
                "or:", font_size=18, bold=True, color=NAVY)

    truth_or = (
        "  A       B     A or B\n"
        "  True    True    True\n"
        "  True    False   True\n"
        "  False   True    True\n"
        "  False   False   False"
    )
    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in truth_or.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = "Courier New"
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(3.9), Inches(5.5), Inches(0.4),
                "not:", font_size=18, bold=True, color=NAVY)

    truth_not = (
        "  A        not A\n"
        "  True     False\n"
        "  False    True"
    )
    txb = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(5.5), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in truth_not.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = "Courier New"
        run.font.color.rgb = GRAY

    add_code_block(s,
                   "# Combining operators\n"
                   "x = 15\n"
                   "print(x > 10 and x < 20)    # True \u2014 between 10 and 20\n"
                   "print(not (x == 15))         # False",
                   Inches(7.0), Inches(1.7), Inches(6.0), Inches(1.8))

    # ── Slide 9: while Loops ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f501 while Loops")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "A  while  loop repeats as long as a condition is  True .",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "count = 1\n"
                   "while count <= 5:\n"
                   "    print(f\"Count: {count}\")\n"
                   "    count += 1   # Don't forget this! (infinite loop otherwise)\n\n"
                   "# Output:\n"
                   "# Count: 1\n"
                   "# Count: 2\n"
                   "# Count: 3\n"
                   "# Count: 4\n"
                   "# Count: 5",
                   Inches(0.8), Inches(2.1), Inches(7.0), Inches(3.8))

    flowchart_while = (
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "\u2502  count <= 5?    \u2502 \u25c4\u2500\u2500\u2500\u2500\u2510\n"
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518      \u2502\n"
        "         \u2502 True          \u2502\n"
        "    \u250c\u2500\u2500\u2500\u2500\u25bc\u2500\u2500\u2500\u2500\u2510          \u2502\n"
        "    \u2502 print   \u251c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "    \u2502 count+= \u2502\n"
        "    \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "         \u2502 False\n"
        "         \u25bc\n"
        "       (done)"
    )
    txb = s.shapes.add_textbox(Inches(8.3), Inches(2.1), Inches(4.5), Inches(3.8))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in flowchart_while.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = "Courier New"
        run.font.color.rgb = NAVY

    # ── Slide 10: break and continue ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\u26d4 break and continue")

    add_code_block(s,
                   "# break \u2014 exit the loop immediately\n"
                   "count = 0\n"
                   'while True:          # "infinite" loop\n'
                   "    count += 1\n"
                   "    if count == 5:\n"
                   "        break        # Stop when count reaches 5\n"
                   "print(f\"Stopped at {count}\")\n\n"
                   "# continue \u2014 skip to next iteration\n"
                   "for i in range(10):\n"
                   "    if i % 2 == 0:\n"
                   "        continue     # Skip even numbers\n"
                   "    print(i)         # Prints: 1, 3, 5, 7, 9",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.1))

    # ── Slide 11: for Loops ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f502 for Loops")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "A  for  loop iterates over a sequence.",
                font_size=20, color=GRAY)

    add_code_block(s,
                   "# range(stop) \u2014 0 to stop-1\n"
                   "for i in range(5):\n"
                   "    print(i)   # 0, 1, 2, 3, 4\n\n"
                   "# range(start, stop)\n"
                   "for i in range(1, 6):\n"
                   "    print(i)   # 1, 2, 3, 4, 5\n\n"
                   "# range(start, stop, step)\n"
                   "for i in range(0, 11, 2):\n"
                   "    print(i)   # 0, 2, 4, 6, 8, 10\n\n"
                   '# Iterate over a string\n'
                   'for letter in "Python":\n'
                   "    print(letter)   # P, y, t, h, o, n",
                   Inches(0.8), Inches(2.1), Inches(11.5), Inches(4.9))

    # ── Slide 12: FizzBuzz ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f3af FizzBuzz \u2014 Classic Coding Challenge")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Print numbers 1 to 30. But:",
                font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.9))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "\u2022  If divisible by 3  \u2192  print \"Fizz\"",
        "\u2022  If divisible by 5  \u2192  print \"Buzz\"",
        "\u2022  If divisible by both  \u2192  print \"FizzBuzz\"",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = item
        run.font.size = Pt(19)
        run.font.color.rgb = GRAY

    add_code_block(s,
                   "for number in range(1, 31):\n"
                   "    if number % 3 == 0 and number % 5 == 0:\n"
                   '        print("FizzBuzz")\n'
                   "    elif number % 3 == 0:\n"
                   '        print("Fizz")\n'
                   "    elif number % 5 == 0:\n"
                   '        print("Buzz")\n'
                   "    else:\n"
                   "        print(number)",
                   Inches(0.8), Inches(3.1), Inches(11.5), Inches(3.0))

    add_textbox(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.45),
                "Output: 1 2 Fizz 4 Buzz Fizz 7 8 Fizz Buzz 11 Fizz 13 14 FizzBuzz ...",
                font_size=17, italic=True, color=TEAL)

    # ── Slide 13: Live Demo: Number Guessing Game ─────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f3ae Live Demo: Number Guessing Game")

    add_code_block(s,
                   "import random\n\n"
                   "secret = random.randint(1, 10)\n"
                   "attempts = 0\n\n"
                   "print(\"I'm thinking of a number between 1 and 10...\")\n\n"
                   "while True:\n"
                   '    guess = int(input("Your guess: "))\n'
                   "    attempts += 1\n\n"
                   "    if guess < secret:\n"
                   '        print("Too low! Try again.")\n'
                   "    elif guess > secret:\n"
                   '        print("Too high! Try again.")\n'
                   "    else:\n"
                   "        print(f\"\U0001f389 Correct! You got it in {attempts} attempts!\")\n"
                   "        break",
                   Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ── Slide 14: Practice Exercises ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f3af Practice Exercises")

    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
                "Try these on your own:", font_size=22, bold=True, color=NAVY)

    exercises = [
        ("Exercise 1", "Write a program that prints all even numbers from 2 to 20."),
        ("Exercise 2",
         'Ask the user for a password. Keep asking until they type "python123".\n'
         '  Then print "Access granted!".'),
        ("Exercise 3",
         "Print a multiplication table for a number the user enters.\n"
         "  e.g. 5 x 1 = 5,  5 x 2 = 10, \u2026  5 x 10 = 50"),
        ("Exercise 4", "Count how many vowels are in a word entered by the user."),
    ]

    y = Inches(2.15)
    for title, body in exercises:
        add_textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.35),
                    title, font_size=20, bold=True, color=TEAL)
        y += Inches(0.35)
        txb = s.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.7))
        tf = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = body
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY
        y += Inches(0.85)

    # ── Slide 15: Assignment 2 ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f4dd Assignment 2: Number Guessing Game")

    add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45),
                "Due: Before Session 3",
                font_size=20, bold=True, color=ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.45),
                "Write a program that:", font_size=20, color=GRAY)

    txb = s.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(11.5), Inches(1.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, t in enumerate([
        "Generates a random number between 1 and 100  (import random)",
        "Asks the user to guess the number repeatedly",
        'Gives hints: "Too high!" or "Too low!"',
        "Counts the number of attempts and prints it when correct",
    ], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {i}. {t}"
        run.font.size = Pt(18)
        run.font.color.rgb = GRAY

    add_textbox(s, Inches(0.8), Inches(4.45), Inches(11.5), Inches(0.4),
                "\u2b50 Bonus:",
                font_size=18, bold=True, color=TEAL)
    txb = s.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.8))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in [
        "Add a maximum of 10 attempts; if exceeded, reveal the answer",
        "Add difficulty levels: Easy (1-50), Medium (1-100), Hard (1-500)",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  \u2022 {b}"
        run.font.size = Pt(17)
        run.font.italic = True
        run.font.color.rgb = TEAL

    add_textbox(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.4),
                "\U0001f4c1 Start with:  starter-code/assignment2_starter.py",
                font_size=17, color=GRAY)
    add_textbox(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4),
                "\U0001f4c4 Full details:  assignments/assignment2_guessing_game.md",
                font_size=17, color=GRAY)

    # ── Slide 16: Session 2 Recap ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, LIGHT)
    add_title_bar(s, "\U0001f389 Session 2 Recap")

    txb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in [
        "Comparison operators:  ==, !=, <, >, <=, >=",
        "if, elif, else \u2014 conditional branching",
        "Logical operators:  and, or, not",
        "while loops \u2014 repeat while condition is True",
        "for loops \u2014 iterate over sequences",
        "range() \u2014 generate number sequences",
        "break \u2014 exit a loop early",
        "continue \u2014 skip to next iteration",
    ]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  \u2705  {item}"
        run.font.size = Pt(21)
        run.font.color.rgb = NAVY

    add_textbox(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "See you in Session 3! \U0001f680",
                font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    return prs


# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os

    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "session2_control_flow.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved {len(prs.slides)} slides \u2192 {out_path}")
